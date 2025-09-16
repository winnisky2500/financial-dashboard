#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Beautify Report Agent
- 输入：经用户确认后的 Markdown 文本 + 可选美化指令（字体、字号、行距、图表配色等）
- 输出：美化后的 HTML（内置 CSS + ECharts 自动渲染）、并上传可下载的 DOCX / PDF（Supabase Storage）
- 鉴权方式、存储桶、可选 LLM 与现有 report_agent 对齐

环境变量（与 report_agent 保持一致）：
  SUPABASE_URL=...
  SUPABASE_SERVICE_ROLE_KEY=...
  REPORTS_BUCKET=reports
  REPORT_AGENT_TOKEN=dev-secret-01

  # 可选：如提供则会用 LLM 对 Markdown 做结构润色（不改动事实内容）
  OPENAI_API_KEY=...
  OPENAI_BASE_URL=https://api.openai.com/v1
  OPENAI_MODEL=gpt-4o-mini

  EXPORT_ENABLED=1  # 1=上传 DOCX/PDF/HTML；0=只返回 HTML 字符串
"""

import os, io, uuid, json, datetime as dt, re, logging, traceback
from typing import Optional, List, Dict, Any

from fastapi import FastAPI, HTTPException, Header, Depends
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from supabase import create_client, Client
import html as html_lib
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import Table, TableStyle
from reportlab.lib import colors

# 可选 LLM
from openai import OpenAI



import matplotlib
matplotlib.use("Agg")           # 无界面环境
from matplotlib import pyplot as plt
# 导出依赖（与 report_agent 保持一致的简单导出策略）
from docx.shared import Pt, Inches
from docx import Document

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader

from matplotlib import rcParams
from urllib.parse import quote

from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image as RLImage,
    ListFlowable, ListItem
)
from reportlab.lib.units import mm
from reportlab.lib.enums import TA_LEFT
from reportlab.lib import colors

import html as html_lib  # 如果上面已有就不要重复导入
# 在现有import之后添加
import re
from typing import Optional, List, Dict, Any, Tuple

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
# ========= DTO =========
class BeautifyStyle(BaseModel):
    # 原有字段保持不变
    font_family: Optional[str] = 'Inter, "Microsoft YaHei", system-ui, -apple-system, Segoe UI, sans-serif'
    heading_font_family: Optional[str] = None
    code_font_family: Optional[str] = 'ui-monospace, SFMono-Regular, Menlo, Monaco, Consolas, "Courier New", monospace'
    base_font_size: Optional[int] = 16
    line_height: Optional[float] = 1.75
    paragraph_spacing_px: Optional[int] = 8
    content_width_px: Optional[int] = 920
    theme: Optional[str] = "light"
    color: Optional[str] = "#111827"
    accent_color: Optional[str] = "#2563eb"
    palette: Optional[List[str]] = None
    
    # 新增用户可配置字段
    template_style: Optional[str] = "modern"          # modern | classic | minimal | business
    show_kpi_cards: Optional[bool] = True             # 是否显示KPI卡片
    show_toc_sidebar: Optional[bool] = True           # 是否显示侧边栏目录
    card_shadow_intensity: Optional[str] = "medium"   # low | medium | high
    border_radius: Optional[int] = 12                 # 圆角大小
    page_background: Optional[str] = "gradient"       # solid | gradient | none
    heading_style: Optional[str] = "modern"           # modern | classic | minimal
    table_style: Optional[str] = "bordered"           # bordered | striped | minimal
    code_theme: Optional[str] = "github"              # github | monokai | tomorrow
    chart_height: Optional[int] = 360                 # 图表高度
    enable_animations: Optional[bool] = True          # 是否启用动画效果
    custom_css: Optional[str] = None                  # 用户自定义CSS


def extract_kpis_from_markdown(md_text: str, show_kpi_cards: bool = True) -> List[Dict[str, str]]:
    """从Markdown中智能提取KPI指标，支持中英文"""
    if not show_kpi_cards:
        return []
    
    kpis = []
    
    # 增强的KPI识别模式 - 支持中英文混合
    patterns = [
        # 中文营收类
        (r'([营收入利润额度规模毛利净利]{2,})\s*[:：]\s*([\d,]+(?:\.\d+)?)\s*([万亿元%]?)', '💰'),
        # 英文营收类
        (r'(Revenue|Income|Profit|Sales|Earnings)\s*[:：]\s*([\d,]+(?:\.\d+)?)\s*([KMBT]?|\$|%)', '💰'),
        
        # 中文增长类
        (r'([增长率变化幅度]{2,}|同比|环比|年化增长|月增长)\s*[:：]\s*([\+\-]?[\d,]+(?:\.\d+)?)\s*([%个点]?)', '📈'),
        # 英文增长类
        (r'(Growth|Increase|Change|YoY|MoM|QoQ)\s*[:：]\s*([\+\-]?[\d,]+(?:\.\d+)?)\s*([%]?)', '📈'),
        
        # 中文占比类
        (r'([市场份额占比比例率转化率]{2,})\s*[:：]\s*([\d,]+(?:\.\d+)?)\s*([%]?)', '📊'),
        # 英文占比类
        (r'(Market Share|Ratio|Rate|Percentage|Share)\s*[:：]\s*([\d,]+(?:\.\d+)?)\s*([%]?)', '📊'),
        
        # 中文数量类
        (r'([用户数量规模总数客户数订单量]{2,})\s*[:：]\s*([\d,]+(?:\.\d+)?)\s*([万个千百]?)', '👥'),
        # 英文数量类
        (r'(Users|Customers|Orders|Count|Volume|Total)\s*[:：]\s*([\d,]+(?:\.\d+)?)\s*([KMBT]?)', '👥'),
        
        # 财务指标
        (r'(ROI|ARPU|客单价|转化率|满意度|CTR|CAC|LTV)\s*[:：]\s*([\d,]+(?:\.\d+)?)\s*([%元$]?)', '🎯'),
        (r'(Conversion|Satisfaction|ROAS|CPC|CPM)\s*[:：]\s*([\d,]+(?:\.\d+)?)\s*([%$]?)', '🎯'),
        
        # 时间相关指标
        (r'([响应时间处理时间平均时间]{3,})\s*[:：]\s*([\d,]+(?:\.\d+)?)\s*([秒分钟小时天]?)', '⏱️'),
        (r'(Response Time|Processing Time|Average Time|Duration)\s*[:：]\s*([\d,]+(?:\.\d+)?)\s*(s|min|h|ms)?', '⏱️'),
    ]
    
    for pattern, icon in patterns:
        matches = re.findall(pattern, md_text, re.IGNORECASE)
        for label, value, unit in matches:
            # 避免重复和无效数据
            if not any(kpi['label'].lower() == label.lower() for kpi in kpis) and value.strip():
                kpis.append({
                    'label': label,
                    'value': f"{value}{unit}",
                    'icon': icon
                })
    
    return kpis[:12]  # 增加到最多显示12个
# --- helpers for table normalization ---
def _normalize_unit(s: str) -> str:
    """将 pct / PCT / Pct 统一成 %，其余不动"""
    return re.sub(r'(?i)pct', '%', s or '')

UNWRAP_MD_TABLE_RE = re.compile(
    r"```(?:markdown|md)\s*\n(\s*\|.+?\|\s*(?:\n\s*\|.+?\|\s*)+)```",
    re.IGNORECASE | re.DOTALL
)
def unwrap_markdown_table_fences(md_text: str) -> str:
    """把 ```markdown 包裹的表格还原成纯表格 Markdown"""
    return UNWRAP_MD_TABLE_RE.sub(lambda m: m.group(1).strip(), md_text or "")

def _delta_polarity(s: str) -> int:
    """判断环比/同比的涨跌：>0 返回 1；<0 返回 -1；否则 0"""
    if not s:
        return 0
    t = s.strip().replace(',', '')
    m = re.search(r'([+\-]?)\s*(\d+(?:\.\d+)?)', t)
    if not m:
        return 0
    sign, num = m.group(1), float(m.group(2))
    if sign == '-': return -1
    if sign == '+': return 1
    return 1 if num > 0 else (-1 if num < 0 else 0)

def inject_kpi_grid(md_text: str, style: Optional[BeautifyStyle] = None) -> str:
    """在Markdown开头注入KPI网格，支持样式配置"""
    if not style or not style.show_kpi_cards:
        return md_text
        
    kpis = extract_kpis_from_markdown(md_text, style.show_kpi_cards)
    
    if not kpis:
        return md_text
    
    # 根据模板样式生成不同的KPI HTML
    template = style.template_style or "modern"
    
    kpi_items = []
    for kpi in kpis:
        if template == "business":
            kpi_items.append(f'''
            <div class="kpi kpi-business">
                <div class="kpi-header">{kpi['icon']}</div>
                <div class="kpi-label">{kpi['label']}</div>
                <div class="kpi-value">{kpi['value']}</div>
            </div>''')
        elif template == "minimal":
            kpi_items.append(f'''
            <div class="kpi kpi-minimal">
                <span class="kpi-label">{kpi['label']}</span>
                <span class="kpi-value">{kpi['value']}</span>
            </div>''')
        else:  # modern (default)
            kpi_items.append(f'''
            <div class="kpi kpi-modern">
                <div class="kpi-label">{kpi['icon']} {kpi['label']}</div>
                <div class="kpi-value">{kpi['value']}</div>
            </div>''')
    
    kpi_html = f'<div class="kpi-grid kpi-{template}">{"".join(kpi_items)}</div>\n\n'
    
    # 在第一个标题前插入KPI
    lines = md_text.split('\n')
    insert_pos = 0
    for i, line in enumerate(lines):
        if line.strip().startswith('#'):
            insert_pos = i
            break
    
    lines.insert(insert_pos, kpi_html)
    return '\n'.join(lines)

_inline_pat = re.compile(r'(\*\*.+?\*\*|\*.+?\*|`.+?`)', re.U)

def _md_inline_to_rl(text: str, bold_font: str, code_font: Optional[str] = None) -> str:
    """
    将 Markdown 内联格式转成 ReportLab 支持的富文本：
    - **粗体** → <font name="{bold_font}">…</font>
    - *斜体*   → <i>…</i>
    - `代码`   → <font name="{code_font或Courier}">…</font>
    其余文本会做 HTML 转义，避免 & / < / > 破坏结构。
    """
    def esc(s: str) -> str:
        return html_lib.escape(s, quote=False)

    parts: list[str] = []
    for tok in re.split(_inline_pat, text or ""):
        if not tok:
            continue
        if tok.startswith("**") and tok.endswith("**"):
            parts.append(f'<font name="{bold_font}">{esc(tok[2:-2])}</font>')
        elif tok.startswith("*") and tok.endswith("*"):
            parts.append(f'<i>{esc(tok[1:-1])}</i>')
        elif tok.startswith("`") and tok.endswith("`"):
            cf = code_font or "Courier"
            parts.append(f'<font name="{cf}">{esc(tok[1:-1])}</font>')
        else:
            parts.append(esc(tok))
    return "".join(parts)

def _parse_font_list(families: Optional[str]) -> list[str]:
    """把 CSS 风格的 font-family 转成名字列表，去掉引号和泛型关键字"""
    if not families:
        return []
    raw = [x.strip().strip("'\"") for x in families.split(",")]
    bad = {"sans-serif", "serif", "monospace", "system-ui", "ui-monospace"}
    return [x for x in raw if x and x not in bad]

def _register_ttf_if_exists(path: str, name: str) -> bool:
    try:
        if os.path.exists(path):
            pdfmetrics.registerFont(TTFont(name, path))
            return True
    except Exception:
        pass
    return False

def _register_family_by_name(name: str) -> tuple[Optional[str], Optional[str]]:
    """
    按常见系统路径尝试注册某个家族的 Regular/Bold。
    返回 (regular_font_name, bold_font_name)；失败返回 (None, None)
    """
    cand: list[tuple[str,str|None,str,str|None]] = []
    # Windows
    cand += [
        ("C:/Windows/Fonts/msyh.ttc", "C:/Windows/Fonts/msyhbd.ttc", "MSYH", "MSYHBD"),           # Microsoft YaHei
        ("C:/Windows/Fonts/simhei.ttf", None, "SimHei", None),                                     # SimHei（无粗体文件，用同名代替）
        ("C:/Windows/Fonts/simsun.ttc", None, "SimSun", None),
        ("C:/Windows/Fonts/Deng.ttf", "C:/Windows/Fonts/Dengb.ttf", "DengXian", "DengXian-Bold"),
    ]
    # Noto / Linux
    cand += [
        ("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
         "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc", "NotoSansCJK", "NotoSansCJK-Bold"),
        ("/usr/share/fonts/truetype/noto/NotoSansSC-Regular.ttf",
         "/usr/share/fonts/truetype/noto/NotoSansSC-Bold.ttf", "NotoSansSC", "NotoSansSC-Bold"),
        ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
         "DejaVuSans", "DejaVuSans-Bold"),
    ]

    # 按名称猜测对应文件
    name_low = name.lower()
    matched = []
    if "yahei" in name_low or "microsoft yahei" in name_low:
        matched = [cand[0]]
    elif "simhei" in name_low:
        matched = [cand[1]]
    elif "simsun" in name_low:
        matched = [cand[2]]
    elif "deng" in name_low:
        matched = [cand[3]]
    elif "noto" in name_low:
        matched = [cand[4], cand[5]]
    elif "dejavu" in name_low:
        matched = [cand[6]]
    else:
        matched = cand  # 不认识就全试一遍

    for reg_path, bold_path, rname, bname in matched:
        ok_r = _register_ttf_if_exists(reg_path, rname)
        ok_b = _register_ttf_if_exists(bold_path, bname) if bold_path else False
        if ok_r:
            return (rname, bname if ok_b else rname)
    return (None, None)

def _resolve_pdf_fonts(style) -> tuple[str, str]:
    """
    根据用户设置解析并注册 PDF 用正文字体/粗体；失败时回退到 CID 字体（STSong）。
    """
    # 1) 用户优先：body/heading 任意一个命中即可
    prefs = _parse_font_list(getattr(style, "font_family", None)) \
          + _parse_font_list(getattr(style, "heading_font_family", None))
    for fam in prefs:
        r, b = _register_family_by_name(fam)
        if r:
            return (r, b or r)

    # 2) 回退：尝试系统常见中文字体
    for fam in ["Microsoft YaHei", "Noto Sans SC", "SimHei", "DejaVu Sans"]:
        r, b = _register_family_by_name(fam)
        if r:
            return (r, b or r)

    # 3) 终极回退：CID 字体（中文可见，但无真粗体）
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        return ("STSong-Light", "STSong-Light")
    except Exception:
        return ("Helvetica", "Helvetica-Bold")

def _apply_docx_font(doc, family: Optional[str]):
    """让 Word 使用用户字体（若未提供则不强制）"""
    fams = _parse_font_list(family)
    name = fams[0] if fams else None
    if not name:
        return
    try:
        from docx.oxml.ns import qn
        style = doc.styles["Normal"]
        style.font.name = name
        r = style._element.get_or_add_rPr()
        rFonts = r.rFonts
        rFonts.set(qn('w:ascii'), name)
        rFonts.set(qn('w:eastAsia'), name)
        rFonts.set(qn('w:hAnsi'), name)
        rFonts.set(qn('w:cs'), name)
    except Exception:
        pass



# ✅ Matplotlib 中文与负号
rcParams['font.sans-serif'] = [
    'Microsoft YaHei', 'SimHei', 'Noto Sans CJK SC',
    'Arial Unicode MS', 'DejaVu Sans', 'sans-serif'
]
rcParams['axes.unicode_minus'] = False


logger = logging.getLogger("beautifyreport")
logger.setLevel(logging.INFO)

# ========= 环境 =========
SUPABASE_URL = os.getenv("SUPABASE_URL") or os.getenv("VITE_SUPABASE_URL", "")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY") or os.getenv("VITE_SUPABASE_SERVICE_ROLE_KEY", "")
REPORTS_BUCKET = os.getenv("REPORTS_BUCKET", "reports")
REPORT_AGENT_TOKEN = os.getenv("REPORT_AGENT_TOKEN", "dev-secret-01")
EXPORT_ENABLED = os.getenv("EXPORT_ENABLED", "1") == "1"

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
OPENAI_BASE_URL = (os.getenv("OPENAI_BASE_URL") or "https://api.openai.com/v1").strip().rstrip("/")
if not OPENAI_BASE_URL.endswith("/v1"):
    OPENAI_BASE_URL += "/v1"
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini")

if not SUPABASE_URL or not SUPABASE_SERVICE_ROLE_KEY:
    raise RuntimeError("缺少 SUPABASE_URL 或 SUPABASE_SERVICE_ROLE_KEY")

sb: Client = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)
llm = OpenAI(api_key=OPENAI_API_KEY or None, base_url=OPENAI_BASE_URL)
def ensure_bucket(bucket: str):
    try:
        # 列出已有桶；不存在则创建为 public
        buckets = sb.storage.list_buckets() or []
        names = [ (b.get("name") if isinstance(b, dict) else getattr(b, "name", None)) for b in buckets ]
        if bucket not in names:
            sb.storage.create_bucket(bucket, {"public": True, "file_size_limit": 104857600})
            logger.info("Created storage bucket: %s", bucket)
    except Exception as e:
        logger.warning("ensure_bucket failed: %s", e)

ensure_bucket(REPORTS_BUCKET)

# ========= FastAPI =========
app = FastAPI(title="Beautify Report Agent", version="1.0.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], allow_credentials=True,
    allow_methods=["*"], allow_headers=["*"],
)

# ========= DTO =========

class BeautifyPayload(BaseModel):
    markdown: str
    language: Optional[str] = "zh"
    instructions: Optional[str] = ""                  # 自然语言美化要求（可选）
    style: Optional[BeautifyStyle] = None

# ========= 安全 =========
def auth_check(authorization: Optional[str] = Header(None)):
    if not REPORT_AGENT_TOKEN:
        return True
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(401, "Missing token")
    if authorization.split(" ", 1)[1] != REPORT_AGENT_TOKEN:
        raise HTTPException(403, "Invalid token")
    return True

# ========= 工具 =========
ECHARTS_BLOCK_RE = re.compile(r"```echarts\s*([\s\S]*?)```", re.MULTILINE)

# ===== 强化 ECharts 配置：legend/tooltip/中文坐标轴 =====
# --- 1) ECharts 选项增强：留白/图例/坐标轴/中文色彩 ---
def _inject_palette(option: dict, palette: Optional[List[str]], theme: str) -> dict:
    if not isinstance(option, dict):
        return option
    opt = json.loads(json.dumps(option, ensure_ascii=False))  # 深拷贝

    # 主题色
    if palette:
        opt["color"] = palette

    # 图例 + 提示框
    opt.setdefault("legend", {})
    if opt["legend"].get("show") is None:
        opt["legend"]["show"] = True
    opt["legend"].setdefault("top", 6)
    opt["legend"].setdefault("left", "center")
    opt.setdefault("tooltip", {"trigger": "axis"})

    # 网格留白：防止图例/坐标轴把图吃掉
    g = opt.setdefault("grid", {})
    g.setdefault("top", 48)
    g.setdefault("left", 56)
    g.setdefault("right", 32)
    g.setdefault("bottom", 48)
    g.setdefault("containLabel", True)

    # 轴样式（暗色/浅色自适应）
    text_color = "#e5e7eb" if (theme or "light").lower() == "dark" else "#374151"
    axis_line = "#4b5563" if (theme or "light").lower() == "dark" else "#d1d5db"

    def _style_axis(ax):
        if not isinstance(ax, dict):
            return
        ax.setdefault("axisLabel", {}).setdefault("color", text_color)
        ax.setdefault("axisLine", {}).setdefault("lineStyle", {})["color"] = axis_line

    # xAxis/yAxis 可能是对象或数组
    if isinstance(opt.get("xAxis"), list):
        for a in opt["xAxis"]:
            _style_axis(a)
            a.setdefault("boundaryGap", isinstance(opt.get("series"), list) and any(
                (isinstance(s, dict) and (s.get("type") == "bar")) for s in (opt.get("series") or [])
            ))
    elif isinstance(opt.get("xAxis"), dict):
        a = opt["xAxis"]
        _style_axis(a)
        a.setdefault("boundaryGap", isinstance(opt.get("series"), list) and any(
            (isinstance(s, dict) and (s.get("type") == "bar")) for s in (opt.get("series") or [])
        ))

    def _style_yaxis(y):
        if isinstance(y, dict):
            _style_axis(y)
            y.setdefault("type", "value")
            y.setdefault("scale", True)  # ★ 高值时不顶到边

    if isinstance(opt.get("yAxis"), list):
        for y in opt["yAxis"]:
            _style_yaxis(y)
    elif isinstance(opt.get("yAxis"), dict):
        _style_yaxis(opt["yAxis"])
    else:
        opt["yAxis"] = {"type": "value", "scale": True}

    # 自动补 series.name，便于 legend 正常显示
    series = opt.get("series") if isinstance(opt.get("series"), list) else []
    names = []
    for i, s in enumerate(series):
        if isinstance(s, dict) and not s.get("name"):
            s["name"] = f"系列{i+1}"
        if isinstance(s, dict):
            names.append(s.get("name") or f"系列{i+1}")
    if names and not opt["legend"].get("data"):
        opt["legend"]["data"] = names

    # 页面全局文本色
    opt.setdefault("textStyle", {}).setdefault("color", text_color)
    return opt



def normalize_echarts_and_extract(md_text: str, palette: Optional[List[str]], theme: str) -> str:
    """把 ```echarts {...}``` 转成 <div class="echarts" data-option="..."></div>，并做 HTML 转义"""
    def _sub(m):
        raw = m.group(1).strip()
        try:
            obj = json.loads(raw)
        except Exception:
            return m.group(0)  # 原样保留
        obj = _inject_palette(obj, palette, theme)
        dataset = json.dumps(obj, ensure_ascii=False)
        safe = html_lib.escape(dataset, quote=True)  # ★ 关键：转义
        # 用双引号包裹属性，避免被单引号截断
        return f'<div class="echarts" data-option="{safe}" style="height:360px;margin:12px 0;"></div>'
    return ECHARTS_BLOCK_RE.sub(_sub, md_text)


def md_to_html_naive(md: str) -> str:
    """轻量 Markdown -> HTML；升级表格解析：跳过分隔行、给环比/同比上色、统一 pct→%"""
    html = md

    # 代码块（除 echarts 已替换外）
    html = re.sub(
        r"```([a-zA-Z0-9_-]+)?\n([\s\S]*?)```",
        lambda m: f'<pre class="code"><code>{m.group(2).replace("<","&lt;").replace(">","&gt;")}</code></pre>',
        html
    )

    # 标题
    html = re.sub(r"^### (.*)$", r'<h3>\1</h3>', html, flags=re.MULTILINE)
    html = re.sub(r"^## (.*)$",  r'<h2>\1</h2>', html, flags=re.MULTILINE)
    html = re.sub(r"^# (.*)$",   r'<h1>\1</h1>', html, flags=re.MULTILINE)

    # 列表
    html = re.sub(r"^\- (.*)$",  r'<li>\1</li>', html, flags=re.MULTILINE)
    html = re.sub(r"^\d+\.\s+(.*)$", r'<li>\1</li>', html, flags=re.MULTILINE)
    html = re.sub(r"((?:<li>.*?</li>\n?)+)", r"<ul>\1</ul>", html)  # 连续 li 包成 ul

    # 粗斜体
    html = re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", html)
    html = re.sub(r"\*(.*?)\*",     r"<em>\1</em>", html)

    # 引用
    html = re.sub(r"^> (.*)$", r'<blockquote>\1</blockquote>', html, flags=re.MULTILINE)

    # ===== 升级版表格解析 =====
    lines = html.splitlines()
    out: list[str] = []
    i = 0
    row_re = re.compile(r'^\s*\|.*\|\s*$')
    sep_cell_re = re.compile(r'^:?-{3,}:?$')  # --- 或 :---: 之类

    def esc(x: str) -> str:
        return html_lib.escape(x, quote=False)

    while i < len(lines):
        if row_re.match(lines[i] or ''):
            block = []
            while i < len(lines) and row_re.match(lines[i] or ''):
                block.append(lines[i])
                i += 1

            rows: list[list[str]] = []
            for ln in block:
                cells = [c.strip() for c in ln.strip()[1:-1].split('|')]
                # 跳过全是 ---/:---: 的分隔行
                if cells and all(sep_cell_re.match(c or '') for c in cells):
                    continue
                rows.append(cells)

            if rows:
                header = rows[0]
                body   = rows[1:]

                # 找出“环比/同比/QoQ/YoY/MoM”列
                delta_idx = set()
                for idx, h in enumerate(header):
                    hl = h.strip().lower()
                    if ('环比' in hl) or ('同比' in hl) or hl in ('qoq', 'yoy', 'mom'):
                        delta_idx.add(idx)

                thead = "<thead><tr>" + "".join(f"<th>{esc(h)}</th>" for h in header) + "</tr></thead>"
                trs = []
                for r in body:
                    tds = []
                    for j, cell in enumerate(r):
                        txt = _normalize_unit(cell)
                        if j in delta_idx:
                            pol = _delta_polarity(txt)
                            cls = "up" if pol > 0 else ("down" if pol < 0 else "zero")
                            tds.append(f'<td class="delta {cls}">{esc(txt)}</td>')
                        else:
                            tds.append(f"<td>{esc(txt)}</td>")
                    trs.append("<tr>" + "".join(tds) + "</tr>")
                tbody = "<tbody>" + "".join(trs) + "</tbody>"
                out.append(f'<table class="table">{thead}{tbody}</table>')
            continue  # 已消费到非表格行

        # 非表格行原样输出
        out.append(lines[i])
        i += 1

    html = "\n".join(out)

    # 换行
    html = html.replace("\n", "<br />")
    return html

def _slug_id(text: str) -> str:
    # 允许中文作为 id，去掉空白与特殊符号；前缀 sec-
    t = re.sub(r"<.*?>", "", text or "")
    t = re.sub(r"\s+", "-", t.strip())
    t = re.sub(r"[^\w\-\u4e00-\u9fff]", "", t)  # 保留中英数/下划线/连字符/中文
    return "sec-" + (t or "section")

def _wrap_section(html: str, sec_id: str, tone: str) -> str:
    # 把指定 id 的 <h2>…</h2> 到下一个 <h2> 前的内容包到卡片里
    pat = re.compile(rf'(<h2 id="{re.escape(sec_id)}"[^>]*>.*?</h2>)(.*?)(?=<h2 id=|$)', re.S)
    def repl(m):
        return f'<section class="card {tone}">' + m.group(1) + m.group(2) + '</section>'
    return pat.sub(repl, html, count=1)

# --- 2) 给 h2/h3 加 id；侧边栏；仅“总起/摘要”做卡片；风险段落粗体变红 ---
def apply_layout_cards_and_toc(inner_html: str, theme: str = "light") -> str:
    html = inner_html

    # 给 h2/h3 加 id，并采集目录
    heads = []
    def add_id(m):
        tag, txt = m.group(1), m.group(2)
        sid = _slug_id(txt)
        heads.append((tag, sid, txt))
        return f'<{tag} id="{sid}">{txt}</{tag}>'
    html = re.sub(r"<(h2|h3)>(.*?)</\1>", add_id, html, flags=re.S)

    # 仅“摘要/总起/总述”包成卡片；风险章节不包卡片，但给 class 方便样式定点强化
    def _wrap_once(sec_id: str, cls: str):
        pat = re.compile(rf'(<h2 id="{re.escape(sec_id)}"[^>]*>.*?</h2>)(.*?)(?=<h2 id=|$)', re.S)
        def repl(m):
            return f'<section class="{cls}">' + m.group(1) + m.group(2) + '</section>'
        return pat.sub(repl, html, count=1)

    for tag, sid, txt in heads:
        if tag != "h2":
            continue
        t = str(txt)
        if any(k in t for k in ["摘要", "总起", "总述"]):
            html = _wrap_once(sid, "card tone-blue")
        elif any(k in t for k in ["风险", "风险点", "风控"]):
            # 只加 class，不加背景
            html = _wrap_once(sid, "risk")

    # 侧边栏目录（只列 h2）
    if heads:
        toc = ["<div class='title'>报告目录</div>"]
        for tag, sid, txt in heads:
            if tag == "h2":
                toc.append(f"<a href='#{sid}'>{html_lib.escape(txt)}</a>")
        toc_html = "<nav>" + "\n".join(toc) + "</nav>"
        html = f"<div class='layout'><aside class='sidebar'>{toc_html}</aside><article>{html}</article></div>"

    return html

# ===== HTML 生成：使用外链 echarts + 内联渲染脚本；支持下载 =====
# --- 3) 页面样式与脚手架：更强的 H1；仅摘要卡片带浅底；风险段落粗体=红字；KPI 宫格可选 ---
def build_html_document(body_inner: str, style: BeautifyStyle) -> str:
    # 解析样式配置 - 添加安全检查
    font = style.font_family or 'Inter, "Microsoft YaHei", system-ui, sans-serif'
    hfont = style.heading_font_family or font
    cfont = style.code_font_family or 'ui-monospace, SFMono-Regular, Menlo, Monaco, Consolas, "Courier New", monospace'
    fs = max(12, style.base_font_size or 16)  # 确保最小字号
    lh = max(1.2, style.line_height or 1.75)  # 确保最小行高
    gap = max(4, style.paragraph_spacing_px or 8)  # 确保最小间距
    width = max(600, style.content_width_px or 920)  # 确保最小宽度
    theme = (style.theme or "light").lower()
    
    # 颜色安全检查 - 确保都是有效的颜色值
    def safe_color(color_val, fallback):
        if not color_val or not color_val.startswith('#'):
            return fallback
        return color_val
    
    color = safe_color(style.color, "#111827")
    accent = safe_color(style.accent_color, "#2563eb")
    
    # 其余代码保持不变...

    
    # 新增样式配置
    template = style.template_style or "modern"
    radius = style.border_radius or 12
    shadow_intensity = style.card_shadow_intensity or "medium"
    page_bg_type = style.page_background or "gradient"
    
    # 根据模板和主题计算颜色
    if theme == "dark":
        bg = "#0b1220"
        if page_bg_type == "gradient":
            page_bg = "linear-gradient(135deg, #1e1b4b 0%, #312e81 50%, #1e1b4b 100%)"
        elif page_bg_type == "solid":
            page_bg = "#1a1a2e"
        else:
            page_bg = bg
        subtle_border = "rgba(255,255,255,.14)"
        subtle_bg = "rgba(255,255,255,.06)"
    else:
        bg = "#ffffff"
        if page_bg_type == "gradient":
            if template == "business":
                page_bg = "linear-gradient(135deg, #667eea 0%, #764ba2 50%, #667eea 100%)"
            elif template == "minimal":
                page_bg = "#f8fafc"
            else:  # modern
                page_bg = "linear-gradient(135deg, #667eea 0%, #764ba2 50%, #667eea 100%)"
        elif page_bg_type == "solid":
            page_bg = "#f1f5f9"
        else:
            page_bg = bg
        subtle_border = "rgba(139, 92, 246, 0.15)"
        subtle_bg = "rgba(37,99,235,.08)"
    
    # 阴影强度
    shadow_map = {
        "low": "0 4px 12px rgba(37, 99, 235, 0.08)" if theme == "light" else "0 4px 12px rgba(0,0,0,.2)",
        "medium": "0 12px 32px rgba(37, 99, 235, 0.15)" if theme == "light" else "0 8px 20px rgba(0,0,0,.35)",
        "high": "0 20px 40px rgba(37, 99, 235, 0.2)" if theme == "light" else "0 12px 32px rgba(0,0,0,.5)"
    }
    card_shadow = shadow_map.get(shadow_intensity, shadow_map["medium"])
    
    muted_text = "#9ca3af" if theme == "dark" else "#6b7280"

    # Web字体导入
    web_fonts = []
    font_families = _parse_font_list(style.font_family)
    for font_name in font_families:
        if 'inter' in font_name.lower():
            web_fonts.append('@import url("https://fonts.googleapis.com/css2?family=Inter:wght@400;600;800&display=swap");')
        elif 'noto' in font_name.lower():
            web_fonts.append('@import url("https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@400;600;800&display=swap");')
    
    web_font_css = '\n'.join(web_fonts)
    
    # 动画CSS
    animations_css = ""
    if style.enable_animations:
        animations_css = """
        @keyframes fadeInUp {
            from { opacity: 0; transform: translateY(20px); }
            to { opacity: 1; transform: translateY(0); }
        }
        @keyframes slideInLeft {
            from { opacity: 0; transform: translateX(-20px); }
            to { opacity: 1; transform: translateX(0); }
        }
        .kpi, .card { animation: fadeInUp 0.6s ease-out; }
        .sidebar { animation: slideInLeft 0.8s ease-out; }
        """

    css = f"""
    {web_font_css}
    {animations_css}
    
    :root {{
      --font: {font};
      --hfont: {hfont};
      --cfont: {cfont};
      --fs: {fs}px;
      --lh: {lh};
      --gap: {gap}px;
      --w: {width}px;
      --fg: {color};
      --bg: {bg};
      --page-bg: {page_bg};
      --accent: {accent};
      --accent2: #8b5cf6;
      --muted: {muted_text};
      --border: {subtle_border};
      --subtle-blue: {subtle_bg};
      --shadow: {card_shadow};
      --radius: {radius}px;
    }}
    
    * {{ box-sizing: border-box; }}
    html,body {{ height: 100%; }}
    body {{
      margin: 0; padding: 24px; color: var(--fg); 
      background: var(--page-bg);
      background-attachment: fixed;
      font-family: var(--font); font-size: var(--fs); line-height: var(--lh);
      -webkit-font-smoothing: antialiased; -moz-osx-font-smoothing: grayscale;
    }}
    
    .container {{ 
      max-width: var(--w); margin: 0 auto; 
      background: var(--bg);
      backdrop-filter: blur(10px);
      border-radius: var(--radius);
      box-shadow: var(--shadow);
      padding: 32px;
      border: 1px solid var(--border);
    }}

    /* ===== 模板样式 ===== */
    
    /* Modern模板 */
    .template-modern h1.hero {{
      font-size: calc(var(--fs) * 2.4);
      font-weight: 800;
      /* 默认字体颜色作为fallback */
      color: var(--accent, #2563eb);
      background: linear-gradient(135deg, var(--accent, #2563eb), var(--accent2, #8b5cf6));
      -webkit-background-clip: text;
      -webkit-text-fill-color: transparent;
      background-clip: text;
      position: relative;
    }}

    /* 不支持background-clip的浏览器fallback */
    @supports not (-webkit-background-clip: text) {{
      .template-modern h1.hero {{
        color: var(--accent, #2563eb) !important;
        background: none !important;
        -webkit-text-fill-color: initial !important;
      }}
    }}

    /* Firefox兼容 */
    @-moz-document url-prefix() {{
      .template-modern h1.hero {{
        color: var(--accent, #2563eb) !important;
        background: none !important;
        -webkit-text-fill-color: initial !important;
      }}
    }}

    .template-modern h1.hero::after {{
      content: ""; position: absolute; left: 0; bottom: 0; height: 4px; width: 120px;
      background: linear-gradient(90deg, var(--accent), var(--accent2));
      border-radius: 8px;
    }}
    
    /* Business模板 */
    .template-business {{
      --accent: #1e40af;
      --accent2: #3b82f6;
    }}
    .template-business h1.hero {{
      font-size: calc(var(--fs) * 2.2);
      font-weight: 700;
      color: var(--accent);
      border-bottom: 3px solid var(--accent);
      padding-bottom: 16px;
    }}
    .template-business .kpi-business {{
      background: linear-gradient(135deg, #1e40af, #3b82f6);
      border-radius: 8px;
    }}
    
    /* Minimal模板 */
    .template-minimal {{
      --radius: 4px;
    }}
    .template-minimal h1.hero {{
      font-size: calc(var(--fs) * 2.0);
      font-weight: 600;
      color: var(--accent);
      border-left: 4px solid var(--accent);
      padding-left: 20px;
    }}
    .template-minimal .kpi-minimal {{
      background: transparent;
      border: 2px solid var(--border);
      border-radius: 4px;
      padding: 12px;
      display: flex;
      justify-content: space-between;
      align-items: center;
    }}

    /* ===== 通用元素样式 ===== */
    
    h1,h2,h3,h4,h5,h6 {{ 
      font-family: var(--hfont); 
      margin: 1.2em 0 .6em; 
      line-height: 1.25; 
    }}
    
    h2 {{ 
      font-size: calc(var(--fs) * 1.6); 
      border-left: 4px solid var(--accent); 
      padding-left: .55em;
      position: relative;
    }}
    h2::before {{
      content: "";
      position: absolute;
      left: -4px;
      top: 0;
      bottom: 0;
      width: 4px;
      background: linear-gradient(180deg, var(--accent), var(--accent2));
    }}
    h3 {{ font-size: calc(var(--fs) * 1.3); color: var(--accent); }}

    p, ul, ol, blockquote, table, pre {{ margin: var(--gap) 0; }}
    ul, ol {{ padding-left: 1.2em; }}
    strong {{ font-weight: 600; color: var(--accent); }}
    em {{ font-style: italic; }}
    a {{ color: var(--accent); text-decoration: none; }}
    a:hover {{ text-decoration: underline; }}

    /* ===== KPI样式增强 ===== */
    
    .kpi-grid {{ 
        display: grid; 
        grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); 
        gap: 16px; 
        margin: 20px 0; 
    }}
    
    .kpi {{ 
        padding: 20px; 
        border: 1px solid var(--border); 
        border-radius: var(--radius); 
        position: relative;
        overflow: hidden;
        transform: translateY(0);
        transition: all 0.3s ease;
    }}
    
    .kpi-modern {{
        background: linear-gradient(135deg, var(--accent), var(--accent2));
        color: white;
    }}
    .kpi-modern:hover {{
        transform: translateY(-4px);
        box-shadow: 0 12px 24px rgba(37, 99, 235, 0.3);
    }}
    .kpi-modern .kpi-label {{ 
        font-size: 13px; 
        color: rgba(255,255,255,0.9); 
        display: flex; 
        align-items: center; 
        gap: 8px; 
        margin-bottom: 8px;
    }}
    .kpi-modern .kpi-value {{ 
        font-size: 24px; 
        font-weight: 700; 
        color: white;
        text-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }}
    
    .kpi-business {{
        background: linear-gradient(135deg, #1e40af, #3b82f6);
        color: white;
    }}
    .kpi-business .kpi-header {{
        font-size: 24px;
        margin-bottom: 8px;
    }}
    .kpi-business .kpi-label {{
        font-size: 12px;
        opacity: 0.9;
        text-transform: uppercase;
        letter-spacing: 0.5px;
    }}
    .kpi-business .kpi-value {{
        font-size: 28px;
        font-weight: 800;
        margin-top: 4px;
    }}
    
    .kpi-minimal {{
        background: transparent;
        border: 2px solid var(--border);
        color: var(--fg);
    }}
    .kpi-minimal .kpi-label {{
        font-size: 14px;
        color: var(--muted);
    }}
    .kpi-minimal .kpi-value {{
        font-size: 20px;
        font-weight: 700;
        color: var(--accent);
    }}

    /* ===== 卡片和布局 ===== */
    
    .card {{
        padding: 20px 24px; 
        border: 1px solid var(--border); 
        border-radius: var(--radius);
        box-shadow: var(--shadow); 
        margin: 16px 0;
        background: transparent;
        position: relative;
        overflow: hidden;
    }}
    .card::before {{
        content: "";
        position: absolute;
        top: 0;
        left: 0;
        right: 0;
        height: 4px;
        background: linear-gradient(90deg, var(--accent), var(--accent2));
    }}
    .card.tone-blue {{ 
        background: var(--subtle-blue); 
        border-color: var(--accent);
    }}

    .risk strong {{ color: #ef4444; }}

    /* ===== 表格样式 ===== */
    
    table.table {{ 
        width: 100%; 
        border-collapse: collapse; 
        border-radius: var(--radius); 
        overflow: hidden;
        box-shadow: var(--shadow);
        margin: 16px 0;
    }}
    table.table th, table.table td {{ 
        border-bottom: 1px solid var(--border); 
        padding: 12px 16px; 
        text-align: left;
    }}
    table.table thead th {{ 
        background: linear-gradient(135deg, var(--accent), var(--accent2)); 
        color: white;
        font-weight: 600;
    }}
    table.table tbody tr:hover {{
        background-color: var(--subtle-blue);
    }}
    table.table td.delta.up {{ color: #10b981; font-weight: 600; }}
    table.table td.delta.down {{ color: #ef4444; font-weight: 600; }}
    table.table td.delta.zero {{ color: var(--muted); }}



    /* ===== 代码块样式 ===== */
    
    pre, code {{ font-family: var(--cfont); }}
    pre.code {{ 
        background: linear-gradient(135deg, rgba(0,0,0,.05), rgba(0,0,0,.02)); 
        padding: 16px; 
        border-radius: var(--radius); 
        overflow: auto;
        border: 1px solid var(--border);
        position: relative;
    }}
    pre.code::before {{
        content: "";
        position: absolute;
        top: 0;
        left: 0;
        width: 4px;
        height: 100%;
        background: var(--accent);
    }}

    /* ===== 侧边栏 ===== */
    
    .layout {{ display: grid; grid-template-columns: 280px 1fr; gap: 32px; }}
    .sidebar {{
        position: sticky; top: 24px; align-self: start; 
        background: var(--subtle-blue);
        border: 1px solid var(--border);
        border-radius: var(--radius);
        padding: 20px;
        max-height: calc(100vh - 100px); 
        overflow: auto;
    }}
    .sidebar nav a {{ 
        display: block; 
        padding: 10px 12px; 
        color: var(--muted); 
        text-decoration: none;
        border-radius: calc(var(--radius) / 2);
        margin: 2px 0;
        transition: all 0.2s ease;
    }}
    .sidebar nav a:hover {{ 
        color: var(--accent); 
        background: rgba(37, 99, 235, 0.1);
    }}
    .sidebar .title {{ 
        font-weight: 600; 
        color: var(--accent); 
        margin-bottom: 12px; 
        font-size: 16px;
    }}

    /* ===== 图表容器 ===== */
    
    .echarts {{
        width: 100%; 
        height: {style.chart_height or 360}px; 
        border: 1px solid var(--border); 
        border-radius: var(--radius); 
        margin: 16px 0;
        background: var(--bg);
        box-shadow: var(--shadow);
    }}

    /* ===== 响应式设计 ===== */
    
    @media (max-width: 1100px) {{ 
        .layout {{ grid-template-columns: 1fr; }} 
        .sidebar {{ position: relative; max-height: none; }} 
        .kpi-grid {{ grid-template-columns: repeat(2, 1fr); }}
    }}
    
    @media (max-width: 768px) {{
        .container {{ padding: 20px; }}
        .kpi-grid {{ grid-template-columns: 1fr; }}
        h1.hero {{ font-size: calc(var(--fs) * 1.8) !important; }}
    }}

    .footer {{ 
        color: var(--muted); 
        font-size: 12px; 
        margin-top: 40px; 
        text-align: center;
        padding-top: 20px;
        border-top: 1px solid var(--border);
    }}
    
    /* 用户自定义CSS */
    {style.custom_css or ''}
    """

    # 应用模板类到body
    template_class = f"template-{template}"
    
    # 其余代码保持不变...
    # 强化 H1：把第一处 <h1> 加 hero class
    body_inner = re.sub(r"<h1>(.*?)</h1>", r'<h1 class="hero">\1</h1>', body_inner, count=1)
    
    # 处理侧边栏显示
    if not style.show_toc_sidebar:
        body_inner = re.sub(r'<div class=[\'"]layout[\'"]><aside class=[\'"]sidebar[\'"]>.*?</aside><article>(.*?)</article></div>', r'\1', body_inner, flags=re.S)

    echarts_script = '<script src="https://cdn.jsdelivr.net/npm/echarts@5/dist/echarts.min.js" defer></script>'
    js = """
    document.addEventListener('DOMContentLoaded', function(){
      if(!window.echarts) return;
      var nodes = document.querySelectorAll('.echarts');
      nodes.forEach(function(div){
        try{
          var raw = div.getAttribute('data-option') || '{}';
          var opt = JSON.parse(raw);
          var chart = echarts.init(div);
          chart.setOption(opt || {});
          setTimeout(function(){ chart.resize(); }, 30);
          window.addEventListener('resize', function(){ chart && chart.resize(); });
        }catch(e){ console.warn('echarts parse failed', e); }
      });
    });
    """

    html = f"""<!DOCTYPE html>
<html lang="zh" class="{template_class}">
<head>
<meta charset="utf-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<meta http-equiv="Content-Security-Policy" content="default-src 'self' https: data: 'unsafe-inline' 'unsafe-eval';">
<title>Beautified Report</title>
<style>{css}</style>
{echarts_script}
</head>
<body class="{template_class}">
  <main class="container">
    {body_inner}
  </main>
<script>{js}</script>
</body>
</html>
"""
    return html






# --- 改后（整段替换） ---
# ===== 上传：确保 contentType 为字符串，并返回可访问链接 =====
def _upload(path: str, content: bytes, content_type: str) -> str:
    # 注意：supabase-py 会把 upsert 放到请求头 x-upsert，必须是 "true"/"false" 字符串
    sb.storage.from_(REPORTS_BUCKET).upload(
        path,
        content,
        {"contentType": str(content_type), "upsert": "true"}
    )
    resp = sb.storage.from_(REPORTS_BUCKET).get_public_url(path)
    if isinstance(resp, dict):
        url = resp.get("publicUrl") or (resp.get("data") or {}).get("publicUrl") \
              or (resp.get("data") or {}).get("public_url")
    else:
        url = resp
    if not isinstance(url, str) or not url:
        raise RuntimeError("get_public_url 返回空")
    return url

def _make_download_url(path: str, public_url: str, filename: str) -> str:
    try:
        resp = sb.storage.from_(REPORTS_BUCKET).create_signed_url(path, 60*60*24, {"download": filename})
        if isinstance(resp, dict):
            signed = resp.get("signed_url") or (resp.get("data") or {}).get("signed_url")
            if signed: return signed
    except Exception as e:
        logger.warning("create_signed_url failed: %s", e)
    # 兜底：把 public_url 当普通字符串处理
    sep = "&" if isinstance(public_url, str) and "?" in public_url else "?"
    return f"{public_url}{sep}download={quote(filename)}"



# ===== 把 ECharts 渲成 PNG（供 DOCX/PDF 使用），保证中文 =====
def _render_chart_png(opt: dict, style: BeautifyStyle) -> bytes:
    """
    用 matplotlib 近似渲染 ECharts（支持 line/bar），并优先使用用户字体。
    """
    try:
        # ★ 用户字体优先
        fams = _parse_font_list(getattr(style, "font_family", None))
        if fams:
            rcParams['font.sans-serif'] = fams + rcParams.get('font.sans-serif', [])
        rcParams['axes.unicode_minus'] = False

        # --- robust x/y extraction (替换原有的 x_data/series 抽取逻辑) ---
        x_data = []
        xa = opt.get("xAxis")
        if isinstance(xa, dict):
            x_data = (xa.get("data") or [])[:]
        elif isinstance(xa, list) and xa:
            xa0 = xa[0]
            if isinstance(xa0, dict):
                x_data = (xa0.get("data") or [])[:]

        # dataset.source 兜底
        if not x_data and isinstance(opt.get("dataset"), dict):
            src = opt["dataset"].get("source") or []
            if src and isinstance(src[0], list):
                # 若首行为表头，取第一列为类目
                head = src[0]
                body = src[1:] if any(isinstance(v, str) for v in head) else src
                x_data = [row[0] for row in body if isinstance(row, list)]
            elif src and isinstance(src[0], dict):
                dim0 = list(src[0].keys())[0]
                x_data = [row.get(dim0) for row in src]

        series = opt.get("series")
        series = series if isinstance(series, list) else ([series] if isinstance(series, dict) else [])

        def _values(arr):
            out = []
            for v in (arr or []):
                if isinstance(v, dict):
                    out.append(v.get("value"))
                else:
                    out.append(v)
            return out

        # 若仍无 x_data，用第一个序列长度生成 1..N
        if not x_data and series and isinstance(series[0], dict):
            n = len(_values(series[0].get("data")))
            x_data = list(range(1, n + 1))

        # 无有效序列直接返回空
        if not series:
            return b""


        palette = opt.get("color") or (style.palette or
                   ["#2563eb","#10b981","#f59e0b","#ef4444","#8b5cf6","#06b6d4"])
        fig = plt.figure(figsize=(7.2, 4.2), dpi=160)
        ax = fig.gca()
        x_idx = list(range(len(x_data)))

        only_bar = all((s.get("type") == "bar") for s in series if isinstance(s, dict))
        for i, s in enumerate(series):
            y = s.get("data") or []
            name = s.get("name") or f"系列{i+1}"
            typ = s.get("type") or ("bar" if only_bar else "line")
            color = palette[i % len(palette)]
            if typ == "bar":
                n = len(series); width = 0.8 / max(1, n)
                ax.bar([t + (i - (n-1)/2)*width for t in x_idx], y, width=width, label=name, color=color)
            else:
                ax.plot(x_idx, y, marker='o', label=name, color=color)

        ax.set_xticks(x_idx)
        ax.set_xticklabels([str(x) for x in x_data], rotation=30, ha='right')
        ax.grid(True, linestyle='--', alpha=0.25)
        if isinstance(opt.get("title"), dict):
            t = opt["title"].get("text") or ""
            if t: ax.set_title(t)
        ax.legend()

        buf = io.BytesIO()
        plt.tight_layout()
        fig.savefig(buf, format="png")
        plt.close(fig)
        return buf.getvalue()
    except Exception:
        return b""

def _iter_md_segments(md_text: str):
    """迭代Markdown文本段落，分离文本和ECharts块"""
    pos = 0
    for m in ECHARTS_BLOCK_RE.finditer(md_text):
        if m.start() > pos:
            yield ("text", md_text[pos:m.start()])
        raw = (m.group(1) or "").strip()
        try:
            opt = json.loads(raw)
        except Exception:
            opt = None
        yield ("echarts", opt)
        pos = m.end()
    if pos < len(md_text):
        yield ("text", md_text[pos:])


def _px_to_pt(px: int) -> Pt:
    return Pt(max(8, (px or 16) * 0.75))
INLINE_RE = re.compile(r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`)')
def _add_md_line(doc: Document, text: str, style: BeautifyStyle, bullet: bool = False):
    if bullet:
        try:
            p = doc.add_paragraph(style="List Bullet")
        except KeyError:
            # 某些环境默认模板里没有 List Bullet，回退为普通段落
            p = doc.add_paragraph()
    else:
        p = doc.add_paragraph()

    if style.paragraph_spacing_px:
        p.paragraph_format.space_after = _px_to_pt(style.paragraph_spacing_px)

    for token in re.split(INLINE_RE, text):
        if not token:
            continue
        if token.startswith("**") and token.endswith("**"):
            run = p.add_run(token[2:-2]); run.bold = True
        elif token.startswith("*") and token.endswith("*"):
            run = p.add_run(token[1:-1]); run.italic = True
        elif token.startswith("`") and token.endswith("`"):
            run = p.add_run(token[1:-1]); run.font.name = (style.code_font_family or "Consolas")
        else:
            p.add_run(token)



def _iter_md_segments(md_text: str):
    pos = 0
    for m in ECHARTS_BLOCK_RE.finditer(md_text):
        if m.start() > pos:
            yield ("text", md_text[pos:m.start()])
        raw = (m.group(1) or "").strip()
        try:
            opt = json.loads(raw)
        except Exception:
            opt = None
        yield ("echarts", opt)
        pos = m.end()
    if pos < len(md_text):
        yield ("text", md_text[pos:])

def export_docx_from_md(md_text: str, style: BeautifyStyle) -> bytes:
    doc = Document()
    normal = doc.styles["Normal"]
    if style.base_font_size:
        normal.font.size = _px_to_pt(style.base_font_size)
    _apply_docx_font(doc, getattr(style, "font_family", None)) 


    for kind, payload in _iter_md_segments(md_text):
        if kind == "text":
            for raw in (payload or "").splitlines():
                line = raw.strip("\r")
                if line.startswith("# "):
                    doc.add_heading(line[2:].strip(), 0)
                elif line.startswith("## "):
                    doc.add_heading(line[3:].strip(), 1)
                elif line.startswith("### "):
                    doc.add_heading(line[4:].strip(), 2)
                elif line.startswith("- "):
                    _add_md_line(doc, line[2:].strip(), style, bullet=True)
                else:
                    _add_md_line(doc, line, style)
        else:
            opt = _inject_palette(payload or {}, style.palette, style.theme)
            png = _render_chart_png(opt, style)
            if png:
                doc.add_picture(io.BytesIO(png), width=Inches(6.2))
            else:
                _add_md_line(doc, "```echarts ...```", style)

    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()
def export_pptx_from_md(md_text: str, style: BeautifyStyle) -> bytes:
    """
    Markdown → PPTX
    规则：
      - H1/H2 开一个新幻灯片（标题放到 slide title）
      - H3/普通段落/列表，作为要点（bullets）放入当前页内容框
      - Markdown 表格渲染为 PPT 表格
      - ```echarts``` 使用已有 _render_chart_png 转成图片插入
    """
    prs = Presentation()
    layout_title = prs.slide_layouts[0]     # Title
    layout_tac   = prs.slide_layouts[1]     # Title and Content
    layout_blank = prs.slide_layouts[6]     # Blank

    def _new_slide(title: str, with_content: bool = True):
        slide = prs.slides.add_slide(layout_tac if with_content else layout_title)
        slide.shapes.title.text = title or ""
        return slide

    # 当前页缓存
    slide = _new_slide("报告综述", True)
    content_box = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    tf = content_box.text_frame if content_box else None
    if tf:
        tf.clear()

    # 复用你已有的分段器（文本 / echarts）
    for kind, payload in _iter_md_segments(md_text):
        if kind == "text":
            lines = (payload or "").splitlines()
            i = 0
            while i < len(lines):
                raw = lines[i].rstrip()
                line = raw.strip()
                i += 1
                if not line:
                    continue

                # 标题：H1/H2 -> 新页
                m = re.match(r'^(#{1,6})\s*(.+)$', line)
                if m:
                    level = len(m.group(1)); text = m.group(2).strip()
                    if level <= 2:
                        slide = _new_slide(text, True)
                        content_box = slide.placeholders[1] if len(slide.placeholders) > 1 else None
                        tf = content_box.text_frame if content_box else None
                        if tf: tf.clear()
                        continue
                    # H3 起作为 bullet
                    line = m.group(2).strip()

                # 列表：- / 1.
                if line.startswith("- ") or line.startswith("* "):
                    p = tf.add_paragraph() if tf and tf.text else (tf.text_frame if tf else None)
                    if tf and not tf.text:
                        tf.text = line[2:].strip()
                    else:
                        p = tf.add_paragraph() if tf else None
                        if p:
                            p.text = line[2:].strip()
                            p.level = 0
                    continue
                if re.match(r'^\d+\.\s+', line):
                    txt = re.sub(r'^\d+\.\s+', '', line)
                    p = tf.add_paragraph() if tf and tf.text else (tf.text_frame if tf else None)
                    if tf and not tf.text:
                        tf.text = txt
                    else:
                        p = tf.add_paragraph() if tf else None
                        if p:
                            p.text = txt
                            p.level = 0
                    continue

                # 表格：连续管道行
                if "|" in line and line.count("|") >= 2:
                    # 收集表格块
                    tbl_lines = [raw]
                    while i < len(lines) and "|" in lines[i] and lines[i].count("|") >= 2:
                        tbl_lines.append(lines[i].rstrip()); i += 1

                    # 解析（与你 PDF 的表格解析保持一致）
                    rows = []
                    sep_cell_re = re.compile(r'^:?-{3,}:?$')
                    for tline in tbl_lines:
                        cells = [c.strip() for c in tline.strip()[1:-1].split("|")]
                        if cells and all(sep_cell_re.match(c or '') for c in cells):
                            continue
                        rows.append(cells)
                    if not rows:
                        continue

                    header, body = rows[0], rows[1:]
                    # 新开一页放表格，避免内容拥挤
                    slide = _new_slide("分期明细", False)
                    slide = prs.slides.add_slide(layout_blank)
                    left, top, width, height = PptxInches(0.5), PptxInches(1.2), PptxInches(9), PptxInches(5)
                    table = slide.shapes.add_table(rows=len(rows), cols=len(header), left=left, top=top, width=width, height=height).table

                    # 表头
                    for c, text in enumerate(header):
                        cell = table.cell(0, c)
                        cell.text = text
                        cell.text_frame.paragraphs[0].font.bold = True
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(37, 99, 235)  # #2563eb

                    # 内容
                    for r, row in enumerate(body, start=1):
                        for c, text in enumerate(row):
                            cell = table.cell(r, c)
                            cell.text = text

                    continue

                # 普通段落 -> bullet
                if tf:
                    if not tf.text:
                        tf.text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                        p.level = 0

        else:  # echarts -> PNG -> 新页插图
            opt = _inject_palette(payload or {}, style.palette, style.theme)
            png = _render_chart_png(opt, style)
            if not png:
                continue
            slide = prs.slides.add_slide(layout_blank)
            left, top = PptxInches(0.6), PptxInches(1.2)
            pic = slide.shapes.add_picture(io.BytesIO(png), left, top, width=PptxInches(9))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

def _render_chart_png(opt: dict, style: BeautifyStyle) -> bytes:
    """
    用 matplotlib 近似渲染 ECharts（支持 line/bar），并优先使用用户字体。
    """
    try:
        # 用户字体优先
        fams = _parse_font_list(getattr(style, "font_family", None))
        if fams:
            rcParams['font.sans-serif'] = fams + rcParams.get('font.sans-serif', [])
        rcParams['axes.unicode_minus'] = False

        # 提取数据
        x_data = []
        xa = opt.get("xAxis")
        if isinstance(xa, dict):
            x_data = (xa.get("data") or [])[:]
        elif isinstance(xa, list) and xa:
            xa0 = xa[0]
            if isinstance(xa0, dict):
                x_data = (xa0.get("data") or [])[:]

        # dataset.source 兜底
        if not x_data and isinstance(opt.get("dataset"), dict):
            src = opt["dataset"].get("source") or []
            if src and isinstance(src[0], list):
                head = src[0]
                body = src[1:] if any(isinstance(v, str) for v in head) else src
                x_data = [row[0] for row in body if isinstance(row, list)]
            elif src and isinstance(src[0], dict):
                dim0 = list(src[0].keys())[0]
                x_data = [row.get(dim0) for row in src]

        series = opt.get("series")
        series = series if isinstance(series, list) else ([series] if isinstance(series, dict) else [])

        def _values(arr):
            out = []
            for v in (arr or []):
                if isinstance(v, dict):
                    out.append(v.get("value"))
                else:
                    out.append(v)
            return out

        # 若仍无 x_data，用第一个序列长度生成 1..N
        if not x_data and series and isinstance(series[0], dict):
            n = len(_values(series[0].get("data")))
            x_data = list(range(1, n + 1))

        # 无有效序列直接返回空
        if not series:
            return b""

        palette = opt.get("color") or (style.palette or
                   ["#2563eb","#10b981","#f59e0b","#ef4444","#8b5cf6","#06b6d4"])
        fig = plt.figure(figsize=(7.2, 4.2), dpi=160)
        ax = fig.gca()
        x_idx = list(range(len(x_data)))

        only_bar = all((s.get("type") == "bar") for s in series if isinstance(s, dict))
        for i, s in enumerate(series):
            y = s.get("data") or []
            name = s.get("name") or f"系列{i+1}"
            typ = s.get("type") or ("bar" if only_bar else "line")
            color = palette[i % len(palette)]
            if typ == "bar":
                n = len(series); width = 0.8 / max(1, n)
                ax.bar([t + (i - (n-1)/2)*width for t in x_idx], y, width=width, label=name, color=color)
            else:
                ax.plot(x_idx, y, marker='o', label=name, color=color)

        ax.set_xticks(x_idx)
        ax.set_xticklabels([str(x) for x in x_data], rotation=30, ha='right')
        ax.grid(True, linestyle='--', alpha=0.25)
        if isinstance(opt.get("title"), dict):
            t = opt["title"].get("text") or ""
            if t: ax.set_title(t)
        ax.legend()

        buf = io.BytesIO()
        plt.tight_layout()
        fig.savefig(buf, format="png")
        plt.close(fig)
        return buf.getvalue()
    except Exception:
        return b""


def export_pdf_from_md(md_text: str, style: BeautifyStyle) -> bytes:
    """增强版PDF导出，支持更好的Markdown解析和样式"""
    base_font, bold_font = _resolve_pdf_fonts(style)
    fs = max(8, (style.base_font_size or 16) * 0.75)
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=A4,
        leftMargin=50, rightMargin=50,
        topMargin=60, bottomMargin=60
    )
    
    # 增强样式
    styles = getSampleStyleSheet()
    
    # 根据用户配置调整样式
    normal_style = styles['Normal']
    normal_style.fontName = base_font
    normal_style.fontSize = fs
    normal_style.leading = fs * (style.line_height or 1.75)
    normal_style.spaceAfter = style.paragraph_spacing_px or 8
    
    # 标题样式增强
    h1_style = styles['Heading1']
    h1_style.fontName = bold_font
    h1_style.fontSize = fs * 1.8
    h1_style.textColor = colors.HexColor(style.accent_color or '#2563eb')
    h1_style.spaceAfter = 16
    h1_style.spaceBefore = 20
    
    h2_style = styles['Heading2']
    h2_style.fontName = bold_font
    h2_style.fontSize = fs * 1.4
    h2_style.textColor = colors.HexColor(style.accent_color or '#2563eb')
    h2_style.spaceAfter = 12
    h2_style.spaceBefore = 16
    h2_style.leftIndent = 0
    h2_style.borderWidth = 0
    h2_style.borderColor = colors.HexColor(style.accent_color or '#2563eb')
    h2_style.borderPadding = (0, 0, 0, 8)
    
    h3_style = styles['Heading3']
    h3_style.fontName = bold_font  
    h3_style.fontSize = fs * 1.2
    h3_style.textColor = colors.HexColor(style.accent_color or '#2563eb')
    
    # 创建彩色背景样式
    highlight_style = ParagraphStyle(
        'Highlight',
        parent=normal_style,
        backColor=colors.HexColor('#f0f7ff'),
        borderColor=colors.HexColor(style.accent_color or '#2563eb'),
        borderWidth=1,
        borderPadding=8,
        borderRadius=4
    )
    
    flow = []
    
    # 添加页眉装饰
    def add_page_decoration(canvas, doc):
        canvas.saveState()
        # 兼容设色（优先带 alpha，失败则降级）
        try:
            from reportlab.lib.colors import Color
            canvas.setFillColor(Color(0.15, 0.39, 0.92, alpha=0.1))
        except Exception:
            canvas.setFillColorRGB(0.15, 0.39, 0.92)

        canvas.rect(0, A4[1]-40, A4[0], 40, fill=1, stroke=0)
        canvas.setFont(base_font, 9)
        canvas.setFillColorRGB(0.5, 0.5, 0.5)
        # 页码：reportlab 推荐用 canvas.getPageNumber()
        try:
            page_no = canvas.getPageNumber()
        except Exception:
            page_no = getattr(doc, "page", 1)
        canvas.drawCentredString(A4[0]/2, 30, f"第 {page_no} 页")

        canvas.restoreState()

    
    # 解析Markdown内容
    for kind, payload in _iter_md_segments(md_text):
        if kind == "text":
            lines = (payload or "").splitlines()
            i = 0
            while i < len(lines):
                line = lines[i].strip()
                
                if not line:
                    i += 1
                    continue
                
                # 增强的标题解析 - 更容错
                header_match = re.match(r'^(#{1,6})\s*(.+)$', line)
                if header_match:
                    level = len(header_match.group(1))
                    title_text = header_match.group(2).strip()
                    
                    if level == 1:
                        flow.append(Paragraph(_md_inline_to_rl(title_text, bold_font), h1_style))
                    elif level == 2:
                        flow.append(Paragraph(_md_inline_to_rl(title_text, bold_font), h2_style))
                    elif level >= 3:
                        flow.append(Paragraph(_md_inline_to_rl(title_text, bold_font), h3_style))
                
                # 处理没有空格的井号标题（容错）
                elif re.match(r'^#{1,6}[^#\s]', line):
                    level = 0
                    for char in line:
                        if char == '#':
                            level += 1
                        else:
                            break
                    title_text = line[level:].strip()
                    
                    if title_text:  # 确保有内容
                        if level == 1:
                            flow.append(Paragraph(_md_inline_to_rl(title_text, bold_font), h1_style))
                        elif level == 2:
                            flow.append(Paragraph(_md_inline_to_rl(title_text, bold_font), h2_style))
                        elif level >= 3:
                            flow.append(Paragraph(_md_inline_to_rl(title_text, bold_font), h3_style))
                
                # 列表处理
                elif line.startswith("- ") or line.startswith("* "):
                    list_items = []
                    while i < len(lines) and (lines[i].strip().startswith("- ") or lines[i].strip().startswith("* ")):
                        item_text = lines[i].strip()[2:].strip()
                        list_items.append(ListItem(Paragraph(_md_inline_to_rl(item_text, bold_font), normal_style)))
                        i += 1
                    flow.append(ListFlowable(list_items, bulletType='bullet'))
                    continue
                
                # 编号列表
                elif re.match(r'^\d+\.\s+', line):
                    list_items = []
                    while i < len(lines) and re.match(r'^\d+\.\s+', lines[i].strip()):
                        item_text = re.sub(r'^\d+\.\s+', '', lines[i].strip())
                        list_items.append(ListItem(Paragraph(_md_inline_to_rl(item_text, bold_font), normal_style)))
                        i += 1
                    flow.append(ListFlowable(list_items, bulletType='1'))
                    continue
                
                # 引用块
                elif line.startswith("> "):
                    quote_text = line[2:].strip()
                    flow.append(Paragraph(_md_inline_to_rl(quote_text, bold_font), highlight_style))
                
                # 表格检测和处理
                elif "|" in line and line.count("|") >= 2:
                    table_lines = [line]
                    j = i + 1
                    while j < len(lines) and "|" in lines[j] and lines[j].count("|") >= 2:
                        table_lines.append(lines[j]); j += 1

                    # 解析表格：跳过分隔行（--- / :---:）
                    rows: list[list[str]] = []
                    sep_cell_re = re.compile(r'^:?-{3,}:?$')
                    for tline in table_lines:
                        cells = [c.strip() for c in tline.strip()[1:-1].split("|")]
                        if cells and all(sep_cell_re.match(c or '') for c in cells):
                            continue
                        rows.append(cells)

                    if rows:
                        header = rows[0]
                        body   = rows[1:]

                        # 统一单位（pct -> %）
                        table_data = [header] + [[_normalize_unit(c) for c in r] for r in body]

                        # 构建表格
                        table = Table(table_data)
                        ts = [
                            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(style.accent_color or '#2563eb')),
                            ('TEXTCOLOR',  (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN',      (0, 0), (-1, -1), 'LEFT'),
                            ('FONTNAME',   (0, 0), (-1, 0), bold_font),
                            ('FONTNAME',   (0, 1), (-1, -1), base_font),
                            ('FONTSIZE',   (0, 0), (-1, -1), fs * 0.9),
                            ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                            ('TOPPADDING',    (0, 0), (-1, -1), 8),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f8fafc')),
                            ('GRID',       (0, 0), (-1, -1), 1, colors.HexColor('#e2e8f0')),
                        ]

                        # 环比/同比列上色
                        delta_idx = [idx for idx, h in enumerate(header)
                                     if ('环比' in h) or ('同比' in h) or h.strip().lower() in ('qoq','yoy','mom')]
                        for r_i, row_vals in enumerate(body, start=1):  # 从第1行（非表头）开始
                            for c_i in delta_idx:
                                val = row_vals[c_i] if c_i < len(row_vals) else ''
                                pol = _delta_polarity(val)
                                if pol > 0:
                                    col = colors.HexColor('#10b981')
                                elif pol < 0:
                                    col = colors.HexColor('#ef4444')
                                else:
                                    col = colors.HexColor('#6b7280')
                                ts.append(('TEXTCOLOR', (c_i, r_i), (c_i, r_i), col))
                                ts.append(('FONTNAME',  (c_i, r_i), (c_i, r_i), bold_font))

                        table.setStyle(TableStyle(ts))
                        flow.append(table)
                        flow.append(Spacer(1, 8))
                        i = j - 1
                    else:
                        flow.append(Paragraph(_md_inline_to_rl(line, bold_font), normal_style))
                
                # 普通段落
                else:
                    if line:  # 非空行才添加
                        flow.append(Paragraph(_md_inline_to_rl(line, bold_font), normal_style))
                
                i += 1
                
        elif kind == "echarts":
            # 图表处理
            opt = _inject_palette(payload or {}, style.palette, style.theme)
            png_data = _render_chart_png(opt, style)
            if png_data:
                try:
                    img = RLImage(io.BytesIO(png_data))
                    img.drawHeight = (style.chart_height or 360) * 0.75  # PDF中缩小一些
                    img.drawWidth = 400
                    flow.append(img)
                    flow.append(Spacer(1, 12))
                except Exception:
                    flow.append(Paragraph("图表渲染失败", normal_style))
            else:
                flow.append(Paragraph("```echarts 图表```", normal_style))
    
    # 构建PDF
    doc.build(flow, onFirstPage=add_page_decoration, onLaterPages=add_page_decoration)
    return buffer.getvalue()



SYSTEM_PROMPT = (
    "你是【报告排版/美化专家】。你的任务是：在不改动事实和数据的前提下，"
    "把我给你的 Markdown 排成一个**结构清晰、带侧边栏目录、卡片化信息块**的 HTML+Markdown 混合文档。"
    "\n\n【必须遵守】\n"
    "1) **严禁改数据/结论**，只做结构与呈现优化；中文用词保持中性客观。\n"
    "2) **禁止** 引入外部 CSS/JS/框架（例如 Tailwind/Bootstrap/FontAwesome/Chart.js 等）；"
    "   只可输出少量语义化 HTML 容器标签（section/div/aside/nav/span）。\n"
    "3) **图表** 必须保留为 ```echarts {…}``` 代码块（不要改成别的库）；标题下方可补一行图注："
    "   <div class='caption'>说明</div>。\n"
    "4) **表格** 使用 Markdown 表格（第一行是表头），不要手写 <table>。"
    "   粗体/斜体/行内代码继续用 Markdown 语法（**…** / *…* / `…`）。\n"
    "5) **字体** 完全遵循用户设置（不要硬编码字体名称/字号/行距）。\n"
    "\n【页面骨架（请照这个结构组织内容）】\n"
    "== 顶部标题区 ==\n"
    "- 顶部 H1：完整报告标题（例如《XX公司年度财务报告》）。\n"
    "- H1 下给 1~3 行关键信息（时间范围/公司名/报告语言等），用 <span class='tag'>…</span> 包装关键字段。\n"
    "\n"
    "== 主体两栏布局 ==\n"
    "<div class='layout'>\n"
    "  <aside class='sidebar'>\n"
    "    <nav>\n"
    "      <div class='title'>报告目录</div>\n"
    "      <!-- 目录项：锚点链接指向正文标题 id -->\n"
    "      <a href='#sec-摘要'>摘要</a>\n"
    "      <a href='#sec-核心指标'>核心指标</a>\n"
    "      <a href='#sec-分项分析'>分项分析</a>\n"
    "      <a href='#sec-风险与建议'>风险与建议</a>\n"
    "    </nav>\n"
    "    <!-- 可选：关键统计/重点关注列表，使用 .tag 小标签或简短说明 -->\n"
    "  </aside>\n"
    "  <article>\n"
    "    \n"
    "    <!-- 1) 摘要卡片：用 tone-blue/green/red 区分类型 -->\n"
    "    <section id='sec-摘要' class='card tone-blue'>\n"
    "      <h2>摘要</h2>\n"
    "      <!-- 将原始摘要要点合并为条理清晰的段落或列表（不改事实） -->\n"
    "      - 关键结论 1\n"
    "      - 关键结论 2\n"
    "    </section>\n"
    "\n"
    "    <!-- 2) KPI 宫格（若原文有核心指标） -->\n"
    "    <section id='sec-核心指标'>\n"
    "      <h2>核心指标</h2>\n"
    "      <div class='kpi-grid'>\n"
    "        <div class='kpi'>\n"
    "          <div class='label'>利润总额</div>\n"
    "          <div class='value'>**123.4 亿元**</div>\n"
    "          <div class='delta'>同比 +x.xx%</div>\n"
    "        </div>\n"
    "        <!-- 其余 KPI ... -->\n"
    "      </div>\n"
    "    </section>\n"
    "\n"
    "    <!-- 3) 分项分析：每一小节都要有 h2/h3（带 id），图表用 ```echarts```，并加图注 -->\n"
    "    <section id='sec-分项分析'>\n"
    "      <h2>分项分析</h2>\n"
    "      <h3 id='sec-盈利能力'>盈利能力</h3>\n"
    "      ```echarts\n"
    "      { \"title\": {\"text\": \"ROA/ROE 趋势\"}, \"legend\": {}, \"xAxis\": {\"data\": [\"2024Q1\",\"2024Q2\"]},\n"
    "        \"yAxis\": {}, \"series\": [ {\"type\":\"line\",\"name\":\"ROA\",\"data\":[0.12,0.13]},\n"
    "                                   {\"type\":\"line\",\"name\":\"ROE\",\"data\":[0.15,0.16]} ] }\n"
    "      ```\n"
    "      <div class='caption'>ROA/ROE 季度变化</div>\n"
    "      \n"
    "      <!-- 如有多组：营运能力、现金流、风险指标等，每组都同样结构 -->\n"
    "    </section>\n"
    "\n"
    "    <!-- 4) 风险与建议：用不同色调卡片突出 -->\n"
    "    <section id='sec-风险与建议' class='card tone-red'>\n"
    "      <h2>风险与建议</h2>\n"
    "      - 主要风险 1（原因/表现）\n"
    "      - 主要风险 2（原因/表现）\n"
    "    </section>\n"
    "\n"
    "  </article>\n"
    "</div>\n"
    "\n【排版细节】\n"
    "- 所有二级/三级标题都要**带 id**（如 #sec-盈利能力），目录锚点链接要准确指向。\n"
    "- 图表标题简洁，legend 必须开启；图注一句话说明图表代表的含义。\n"
    "- 如果原文含有“重点结论/异常项/提示”，优先用 <section class='card tone-red|green|blue'> 表示；\n"
    "- 如果没有侧边栏所需的信息，也请**至少**输出一个包含二级标题链接的目录。\n"
    "- 不要输出 <html> / <head> / <body> 外层骨架，保持为 Markdown + 少量容器 HTML。\n"
)



@app.get("/health")
def health():
    return {"ok": True, "time": dt.datetime.utcnow().isoformat()}

@app.post("/beautify/run")
def beautify_run(payload: BeautifyPayload, _=Depends(auth_check)):
    try:
        style = payload.style or BeautifyStyle()
        md = payload.markdown or ""
        if not md.strip():
            raise HTTPException(400, "空的 markdown")

        logger.info(f"收到美化请求，样式配置: template={style.template_style}, theme={style.theme}")

        # 1) （可选）LLM轻度排版润色——不改事实
        improved_md = md
        if OPENAI_API_KEY and (payload.instructions or "").strip():
            try:
                prompt = f"""
请对以下Markdown报告进行结构和排版优化，要求：
1. 不修改任何数据、数字、事实内容
2. 优化段落结构和层次
3. 增强可读性和专业性
4. 保持原有的图表代码块不变
5. 根据用户要求进行调整：{payload.instructions}

原始内容：
{md}

请返回优化后的Markdown内容：
"""
                resp = llm.chat.completions.create(
                    model=OPENAI_MODEL,
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.3
                )
                improved_md = resp.choices[0].message.content or md
                logger.info("LLM结构优化完成")
            except Exception as e:
                logger.warning(f"LLM优化失败，使用原始内容: {e}")
                improved_md = md

        # 2) HTML：把 ```echarts``` 替换成 <div class="echarts" ...>，再去掉误包裹的 markdown 围栏
        md_with_charts = normalize_echarts_and_extract(improved_md, style.palette, style.theme)
        md_with_charts = unwrap_markdown_table_fences(md_with_charts)   # ★ 新增一行

        # 3) 注入KPI网格（支持用户配置）
        md_with_kpis = inject_kpi_grid(md_with_charts, style)
        
        # 4) 转换为HTML并应用布局
        body_html = md_to_html_naive(md_with_kpis)
        body_html = apply_layout_cards_and_toc(body_html, style.theme or "light")
        html_doc = build_html_document(body_html, style)

        # 5) 导出/上传
                # 5) 导出/上传（逐项 try，互不影响）
        job_id = str(uuid.uuid4())
        timestamp = dt.datetime.utcnow().strftime("%Y%m%d_%H%M%S")

        result = {
            "job_id": job_id,
            "timestamp": timestamp,
            "style_applied": {
                "template": style.template_style,
                "theme": style.theme,
                "font_family": style.font_family,
                "accent_color": style.accent_color,
                "show_kpi_cards": style.show_kpi_cards,
                "show_toc_sidebar": style.show_toc_sidebar,
            }
        }

        if not EXPORT_ENABLED:
            result["html"] = html_doc
            logger.info("仅返回HTML内容（导出已禁用）")
            return result

        # === HTML ===
        try:
            html_path = f"beautified/{job_id}/report_{timestamp}.html"
            html_bytes = html_doc.encode("utf-8")
            html_pub = _upload(html_path, html_bytes, "text/html")
            result["html_url"] = _make_download_url(html_path, html_pub, f"report_{timestamp}.html")
        except Exception as e:
            logger.error("HTML 导出失败: %s", e)
            # 兜底直接放内联 HTML，前端至少能“下载 HTML”
            result["html"] = html_doc
            result["html_error"] = str(e)

        # === DOCX ===
        try:
            docx_bytes = export_docx_from_md(improved_md, style)
            docx_path = f"beautified/{job_id}/report_{timestamp}.docx"
            docx_pub = _upload(docx_path, docx_bytes,
                               "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            result["docx_url"] = _make_download_url(docx_path, docx_pub, f"report_{timestamp}.docx")
        except Exception as e:
            logger.error("DOCX 导出失败: %s", e)
            result["docx_error"] = str(e)

        # === PDF ===
        try:
            pdf_bytes = export_pdf_from_md(improved_md, style)
            pdf_path = f"beautified/{job_id}/report_{timestamp}.pdf"
            pdf_pub = _upload(pdf_path, pdf_bytes, "application/pdf")
            result["pdf_url"] = _make_download_url(pdf_path, pdf_pub, f"report_{timestamp}.pdf")
        except Exception as e:
            logger.error("PDF 导出失败: %s", e)
            result["pdf_error"] = str(e)

        # === PPTX（新增）===
        try:
            pptx_bytes = export_pptx_from_md(improved_md, style)
            pptx_path = f"beautified/{job_id}/report_{timestamp}.pptx"
            pptx_pub = _upload(pptx_path, pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            result["pptx_url"] = _make_download_url(pptx_path, pptx_pub, f"report_{timestamp}.pptx")
        except Exception as e:
            logger.error("PPTX 导出失败: %s", e)
            result["pptx_error"] = str(e)
            
        return result


    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"处理请求失败: {traceback.format_exc()}")
        raise HTTPException(500, f"处理失败: {str(e)}")

